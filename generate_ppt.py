import os
import re
import sys
from dataclasses import dataclass
from typing import Iterable, Optional


def ensure_stdout_utf8() -> None:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass


@dataclass(frozen=True)
class RunStyle:
    font_name: Optional[str]
    font_size: Optional[int]  # EMU via python-pptx Length is awkward; store raw value
    bold: Optional[bool]
    italic: Optional[bool]
    underline: Optional[bool]
    rgb: Optional[tuple]  # (r,g,b)


def capture_run_style(run) -> RunStyle:
    font = run.font
    rgb = None
    try:
        if font.color and font.color.rgb:
            rgb = (font.color.rgb[0], font.color.rgb[1], font.color.rgb[2])
    except Exception:
        rgb = None

    size = None
    try:
        if font.size is not None:
            size = int(font.size)
    except Exception:
        size = None

    return RunStyle(
        font_name=font.name,
        font_size=size,
        bold=font.bold,
        italic=font.italic,
        underline=font.underline,
        rgb=rgb,
    )


def apply_run_style(run, style: RunStyle) -> None:
    font = run.font
    if style.font_name is not None:
        font.name = style.font_name
    if style.font_size is not None:
        # python-pptx expects a Length; it also accepts int EMU
        font.size = style.font_size
    if style.bold is not None:
        font.bold = style.bold
    if style.italic is not None:
        font.italic = style.italic
    if style.underline is not None:
        font.underline = style.underline
    if style.rgb is not None:
        from pptx.dml.color import RGBColor

        font.color.rgb = RGBColor(*style.rgb)


def normalize(s: str) -> str:
    return " ".join((s or "").strip().split())


def iter_text_shapes(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
            yield shape


def find_shape_by_text(slide, predicate) -> Optional[object]:
    for shape in iter_text_shapes(slide):
        t = normalize(shape.text_frame.text)
        if t and predicate(t):
            return shape
    return None


def find_main_body_shape(slide) -> Optional[object]:
    """Pick the likely 'body' text shape (largest non-title content)."""
    candidates = []
    for shape in iter_text_shapes(slide):
        tf = shape.text_frame
        t = normalize(tf.text)
        if not t:
            continue
        # Exclude slide number-only shapes and obvious headings
        if re.fullmatch(r"\d+", t):
            continue
        if t.upper() in {
            "IDEA TITLE",
            "PROBLEM STATEMENT",
            "PROPOSED SOLUTIONS",
            "TECHNICAL APPROACH",
            "FEASIBILITY & SCALABILITY",
            "TEAM MEMBERS",
        }:
            continue
        # Exclude the big title on slide 2 ("IDEA TITLE" slide) by looking for keywords
        candidates.append(shape)

    if not candidates:
        return None

    def area(sh):
        try:
            return int(sh.width) * int(sh.height)
        except Exception:
            return 0

    candidates.sort(key=area, reverse=True)
    return candidates[0]


def replace_exact_run_text(shape, replacements: dict) -> int:
    """Replace run.text exactly when it matches keys in replacements."""
    tf = shape.text_frame
    changed = 0
    for para in tf.paragraphs:
        for run in para.runs:
            if run.text in replacements:
                run.text = replacements[run.text]
                changed += 1
    return changed


def set_bullets_preserve_style(text_frame, lines: list[str]) -> None:
    """Clear text frame, add paragraphs, preserve style from first available run."""
    style = None
    try:
        for p in text_frame.paragraphs:
            if p.runs:
                style = capture_run_style(p.runs[0])
                break
    except Exception:
        style = None

    # Clear
    text_frame.clear()

    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        # Default bullets are usually on; preserve template behavior by not forcing bullet property.
        # We *do* keep indentation level at 0 for all lines unless the template had others.
        p.text = ""
        r = p.add_run()
        r.text = line
        if style is not None:
            apply_run_style(r, style)


def set_paragraph_lines_preserve_style(text_frame, lines: list[str]) -> None:
    """Like bullets, but uses line breaks as separate paragraphs (no explicit bullet toggles)."""
    set_bullets_preserve_style(text_frame, lines)


def delete_slide(prs, slide_index: int) -> None:
    """Delete slide by index (0-based) using underlying XML relationship removal."""
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001
    sldId_elements = list(slide_id_list)
    sldId = sldId_elements[slide_index]
    # python-pptx stores the relationship id on the element using the 'rId' property
    # (namespaced attribute). Access via attribute when available.
    rId = getattr(sldId, "rId", None)
    if not rId:
        # Fallback to raw namespaced attribute
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
    slide_id_list.remove(sldId)
    if rId:
        prs.part.drop_rel(rId)


def main() -> int:
    ensure_stdout_utf8()
    try:
        from pptx import Presentation
    except Exception:
        print("python-pptx is not installed. Install it with: pip install python-pptx")
        return 2

    template_path = r"C:\Users\MELVIN\Downloads\69c4c39c2d8a0_OSSOMEHACKS_PPTtemplate1 (1).pptx"
    output_path = os.path.join(os.path.dirname(template_path), "neuro_shift_presentation.pptx")

    prs = Presentation(template_path)
    if len(prs.slides) < 10:
        raise RuntimeError(f"Expected 10 slides in template, found {len(prs.slides)}")

    # ---------- Slide 2 (index 1): IDEA TITLE ----------
    slide2 = prs.slides[1]
    # Replace the big title (likely "IDEA TITLE") with "NEURO-SHIFT"
    title_shape = find_shape_by_text(slide2, lambda t: t.upper() == "IDEA TITLE")
    if title_shape is None:
        # fallback: the largest text shape on the slide that isn't slide number
        title_shape = find_main_body_shape(slide2)
    if title_shape is not None:
        # Replace runs safely if possible; else overwrite the text frame with same style
        if replace_exact_run_text(title_shape, {"IDEA TITLE": "NEURO-SHIFT"}) == 0:
            set_paragraph_lines_preserve_style(title_shape.text_frame, ["NEURO-SHIFT"])

    # Fill Team Name / Team Lead / Track line (single shape in template)
    meta_shape = find_shape_by_text(slide2, lambda t: "TEAM NAME:" in t.upper() and "TEAM LEAD:" in t.upper())
    if meta_shape is not None:
        meta_lines = [
            "Team Name: Puranaanooru",
            "Team Lead: Melvin Daniel Gilton Rajasekar",
            "Track: Open Innovation — AI/ML",
        ]
        set_paragraph_lines_preserve_style(meta_shape.text_frame, meta_lines)

    # ---------- Slide 3 (index 2): PROBLEM STATEMENT ----------
    slide3 = prs.slides[2]
    body3 = find_main_body_shape(slide3)
    if body3 is not None:
        bullets = [
            "1 in 4 people globally live with a disability — WHO 2023",
            "For ALS, spinal cord injury, or severe motor impairment patients, controlling lights, fans, and doors requires a caregiver for every single action",
            "Existing solutions (eye-gaze trackers, voice assistants) cost ₹50,000+ and require internet connectivity",
            "No affordable, offline, gesture-based assistive home control exists for low-resource settings in India",
        ]
        set_bullets_preserve_style(body3.text_frame, bullets)

    # ---------- Slide 4 (index 3): PROPOSED SOLUTIONS (Briefly explain) ----------
    slide4 = prs.slides[3]
    body4 = find_main_body_shape(slide4)
    if body4 is not None:
        bullets = [
            "Neuro-Shift reads facial muscle activity (surface EMG) — tiny electrical signals from eyebrow raises and jaw clenches",
            "Converts gestures in real-time into smart home commands: lights, fan, TV, door lock",
            "Runs fully OFFLINE — Arduino + Laptop + ESP32, no cloud, no subscription",
            "Hardware cost: under ₹800 total (Arduino Uno + AD8232 + ESP32 + relay module + electrodes)",
            "Self-calibrates to each individual user in 30 seconds",
            "Designed for differently-abled users excluded from mainstream assistive tech due to cost",
        ]
        set_bullets_preserve_style(body4.text_frame, bullets)

    # ---------- Slide 5 (index 4): PROPOSED SOLUTIONS (How it solves) ----------
    slide5 = prs.slides[4]
    body5 = find_main_body_shape(slide5)
    if body5 is not None:
        bullets = [
            "CAPTURE: AD8232 EMG module reads muscle signals via skin electrodes at 500Hz via Arduino Uno",
            "PROCESS: Laptop receives raw signal over USB serial (pyserial) → 200ms sliding window → extract RMS, MAV, Zero-Crossing Rate, Waveform Length features",
            "CLASSIFY: RandomForest ML model (trained on real user EMG data) → predicts gesture class in <300ms",
            "SAFETY GATE: Confidence below 80% = no action. 1-second cooldown prevents false triggers. Manual override always available",
            "DISPATCH: Flask backend sends HTTP POST to ESP32 over WiFi → ESP32 triggers GPIO relay → controls physical IoT devices",
            "CALIBRATE: 30-sec onboarding records user's own gestures. Model re-fits on personal data — not a generic baseline",
        ]
        set_bullets_preserve_style(body5.text_frame, bullets)

    # ---------- Slide 6 (index 5): TECHNICAL APPROACH (Tech Stack) ----------
    slide6 = prs.slides[5]
    body6 = find_main_body_shape(slide6)
    if body6 is not None:
        lines = [
            "Hardware Layer:",
            "• Arduino Uno + AD8232 EMG module",
            "• Skin surface electrodes (3-lead)",
            "• Laptop / PC — ML inference + Flask backend",
            "• ESP32 — WiFi command receiver + GPIO relay controller",
            "• Relay module (~₹80) → controls smart home devices",
            "",
            "Software Layer:",
            "• Python: NumPy, SciPy (signal processing), scikit-learn (RandomForest classifier)",
            "• pyserial — Arduino to Laptop serial communication over USB",
            "• Flask REST API — gesture command dispatcher",
            "• HTTP POST — Laptop to ESP32 over shared WiFi / hotspot",
            "• emg_rf_model.pkl — trained on real recorded EMG sessions per user",
            "",
            "IoT Outputs:",
            "• Light (on/off), Fan (speed), TV (IR blaster), Door Lock (relay)",
        ]
        set_paragraph_lines_preserve_style(body6.text_frame, lines)

    # ---------- Slide 7 (index 6): TECHNICAL APPROACH (Flowchart/Architecture) ----------
    slide7 = prs.slides[6]
    body7 = find_main_body_shape(slide7)
    if body7 is not None:
        lines = [
            "SKIN ELECTRODES",
            "      ↓",
            "AD8232 Module (amplify + bandpass filter 20–450Hz)",
            "      ↓",
            "Arduino Uno — 500Hz ADC sampling",
            "      ↓ USB Serial (pyserial)",
            "LAPTOP — Python backend",
            "      ↓",
            "Feature Extraction: 200ms window → [RMS, MAV, ZCR, WL]",
            "      ↓",
            "RandomForest Classifier → class probabilities",
            "      ↓",
            "Confidence Gate (≥80%?) — YES → Cooldown Check → HTTP POST",
            "                          — NO  → No Action (idle)",
            "      ↓ WiFi (shared hotspot)",
            "ESP32 — receives JSON command {\"command\": \"light_on\"}",
            "      ↓",
            "GPIO Relay → [ Light | Fan | TV IR Blaster | Door Lock ]",
            "      ↓",
            "LED Feedback to user + 1-sec cooldown reset",
            "",
            "Fallback: Manual override button on ESP32 → direct relay control (bypasses ML entirely)",
        ]
        set_paragraph_lines_preserve_style(body7.text_frame, lines)

    # ---------- Slide 8 (index 7): FEASIBILITY & SCALABILITY ----------
    slide8 = prs.slides[7]
    body8 = find_main_body_shape(slide8)
    if body8 is not None:
        lines = [
            "Feasibility — Works Today:",
            "• Full prototype built: Arduino + AD8232 + Laptop + ESP32 + relay",
            "• 91% per-class accuracy across 3 gesture classes (eyebrow raise, jaw clench, rest)",
            "• End-to-end latency: <300ms (imperceptible to user)",
            "• Laptop + ESP32 on a shared mobile hotspot — no router needed",
            "• 30-sec calibration usable by non-technical caregivers",
            "• Total hardware cost: under ₹800",
            "",
            "Scalability — Growth Path:",
            "• Gesture vocabulary: 2 classes today → 6+ with more training data",
            "• Replace laptop with Raspberry Pi for a fully embedded wearable version",
            "• Multi-room expansion via MQTT broker — software update only",
            "• Electrode form factor → comfortable behind-ear wearable",
            "• Institutional buyers: rehab centres, elderly care homes, govt assistive tech programs",
            "• WHO: 2.5 billion people will need assistive products by 2030",
        ]
        set_paragraph_lines_preserve_style(body8.text_frame, lines)

    # ---------- Slide 9 (index 8): TEAM MEMBERS ----------
    slide9 = prs.slides[8]
    body9 = find_main_body_shape(slide9)
    if body9 is not None:
        lines = [
            "Team Lead: [Your Name] | Email: [email@srmist.edu.in] | Phone: [+91-XXXXXXXXXX]",
            "Member 2: [Name] | Email: [email@srmist.edu.in] | Phone: [+91-XXXXXXXXXX]",
            "Member 3: [Name] | Email: [email@srmist.edu.in] | Phone: [+91-XXXXXXXXXX]",
            "Member 4: [Optional] | Email: [email] | Phone: [phone]",
        ]
        set_paragraph_lines_preserve_style(body9.text_frame, lines)

    # ---------- Delete slide 10 (index 9) ----------
    delete_slide(prs, 9)

    prs.save(output_path)

    # ---------- Validate ----------
    out = Presentation(output_path)
    ok = True
    if len(out.slides) != 9:
        ok = False
        print(f"ERROR: expected 9 slides, got {len(out.slides)}")

    required = [
        "NEURO-SHIFT",
        "Puranaanooru",
        "Open Innovation — AI/ML",
        "91% per-class accuracy",
    ]
    whole_text = []
    for s in out.slides:
        for sh in s.shapes:
            if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                whole_text.append(sh.text_frame.text)
    blob = "\n".join(whole_text)
    for r in required:
        if r not in blob:
            ok = False
            print(f"ERROR: missing required text: {r!r}")

    print(f"Saved: {output_path}")
    print(f"Slide count: {len(out.slides)}")
    print("Validation:", "OK" if ok else "FAILED")
    return 0 if ok else 1


if __name__ == "__main__":
    raise SystemExit(main())

